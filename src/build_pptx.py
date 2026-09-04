#!/usr/bin/env python3
"""Собирает takt2-gamedev.pptx в стиле деки: кремовая бумага, Onest, один акцент.

    python3 src/build_pptx.py            # → takt2-gamedev.pptx в корне репозитория

Иконки берутся из $CLAUDE_JOB_DIR/tmp/pptx/icons или из src/pptx-icons (см. ICONS).
"""
import os, re, json, pathlib, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = pathlib.Path(__file__).resolve().parent.parent
ICONS = pathlib.Path(os.environ.get("PPTX_ICONS", ROOT / "src" / "pptx-icons"))
FONT = "Onest"

def C(h): return RGBColor.from_string(h.lstrip("#"))
PAPER, PAPER2, TILE = C("F1ECE3"), C("E7E1D5"), C("FFFFFF")
INK, MUTED, HAIR = C("16140F"), C("6B6558"), C("D3CCBE")
JET, FLAME, FLAME_INK, SAGE = C("14130F"), C("F26A1B"), C("B34405"), C("DCE2D9")
PAPER_ON_DARK = C("F1ECE3")

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = 13.333, 7.5
PAD = 0.65
CNT = [0]

# ── штриховые иконки → PNG через rsvg-convert ──────────────────────────────
UI = {
 "key":   '<circle cx="8" cy="8" r="4.6"/><path d="M11.3 11.3 L20.5 20.5"/><path d="M17.2 17.2l2.4 2.4"/><path d="M14.4 14.4l2.4 2.4"/>',
 "nomoney":'<rect x="2.5" y="6" width="19" height="12" rx="2"/><path d="M2.5 21.5 L21.5 2.5"/>',
 "compare":'<rect x="2.5" y="5" width="9" height="6" rx="1.4"/><rect x="2.5" y="14" width="15" height="6" rx="1.4"/><path d="M21.5 3.5v17"/>',
 "doc":   '<path d="M6 2.5h8l4 4v15H6z"/><path d="M14 2.5v4h4"/><path d="M9 12h6"/><path d="M9 16h6"/>',
 "gate":  '<path d="M3 21V9l9-6 9 6v12"/><path d="M9 21v-7h6v7"/><path d="M3 21h18"/>',
 "layers":'<path d="M12 3 2.5 8 12 13l9.5-5z"/><path d="M2.5 12.5 12 17.5l9.5-5"/><path d="M2.5 17 12 22l9.5-5"/>',
 "shield":'<path d="M12 2.5 4 6v6c0 5 3.4 8.4 8 9.5 4.6-1.1 8-4.5 8-9.5V6z"/><path d="M9 12l2 2 4-4"/>',
 "eye":   '<path d="M1.5 12S5.5 5 12 5s10.5 7 10.5 7-4 7-10.5 7S1.5 12 1.5 12z"/><circle cx="12" cy="12" r="3"/>',
 "clock": '<circle cx="12" cy="12" r="9"/><path d="M12 7v5l3.5 2"/>',
 "no":    '<circle cx="12" cy="12" r="9"/><path d="M5.6 5.6l12.8 12.8"/>',
 "play":  '<rect x="2.5" y="4.5" width="19" height="15" rx="2.5"/><path d="M10 9v6l5-3z"/>',
 "chain": '<path d="M10 14a4.5 4.5 0 0 0 6.4 0l2.6-2.6a4.5 4.5 0 0 0-6.4-6.4L11.2 6.4"/><path d="M14 10a4.5 4.5 0 0 0-6.4 0L5 12.6a4.5 4.5 0 0 0 6.4 6.4l1.4-1.4"/>',
 "ruler": '<path d="M3 17 17 3l4 4L7 21z"/><path d="M8 12l2 2"/><path d="M11 9l2 2"/><path d="M14 6l2 2"/>',
}
def ui_png(name, color="#F26A1B"):
    p = ICONS / f"ui_{name}_{color.lstrip('#')}.png"
    if not p.exists():
        svg = ('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" fill="none" stroke="%s" '
               'stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round">%s</svg>' % (color, UI[name]))
        s = ICONS / f"ui_{name}.svg"; s.write_text(svg)
        subprocess.run(["rsvg-convert", "-w", "384", "-h", "384", "-o", str(p), str(s)], check=True)
    return str(p)

def ico(slug, on_dark=False):
    p = ICONS / (f"{slug}_paper.png" if on_dark and (ICONS / f"{slug}_paper.png").exists() else f"{slug}.png")
    return str(p)

# ── примитивы ──────────────────────────────────────────────────────────────
def new_slide(bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    CNT[0] += 1
    return s

def rect(s, x, y, w, h, fill, radius=0.13, line=None, line_w=1.5):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if radius: shp.adjustments[0] = min(0.5, radius / min(w, h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def oval(s, x, y, d, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line=1.25, caps=False, spacing=None):
    """runs: строка или список (текст, {bold,color,size,url})."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line
        if spacing: p.space_after = Pt(spacing)
        items = para if isinstance(para, list) else [para]
        for it in items:
            t, o = (it, {}) if isinstance(it, str) else it
            r = p.add_run(); r.text = t.upper() if caps else t
            f = r.font; f.name = FONT; f.size = Pt(o.get("size", size)); f.bold = o.get("bold", bold)
            f.color.rgb = o.get("color", color)
            if o.get("url"): r.hyperlink.address = o["url"]
    return tb

def pic(s, path, x, y, w=None, h=None):
    """Вписывает картинку в бокс w×h по пропорции, центрируя."""
    iw, ih = Image.open(path).size
    if w and h:
        k = min(w / iw, h / ih); pw, ph = iw * k, ih * k
        return s.shapes.add_picture(path, Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))
    return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w) if w else None, Inches(h) if h else None)

def top(s, pill, dark=False, hot=False):
    col = PAPER_ON_DARK if dark else INK
    border = C("6E6A5E") if dark else (C("8A4A28") if hot else INK)
    w = 0.115 * len(pill) + 0.5
    rect(s, PAD, 0.5, w, 0.32, None, radius=0.16, line=border, line_w=1.5)
    text(s, PAD, 0.5, w, 0.32, pill, size=9, bold=True, color=col, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, caps=True)
    text(s, W - PAD - 1.2, 0.5, 1.2, 0.32, f"{CNT[0]:02d} / {{N}}", size=10, bold=True,
         color=(C("9A968C") if dark else (C("8A4A28") if hot else MUTED)), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def h2(s, t, y=1.05, color=INK, size=36):
    text(s, PAD, y, W - 2 * PAD, 1.25, t, size=size, bold=True, color=color, line=1.05)

def note(s, t, y=6.75, color=MUTED):
    text(s, PAD, y, W - 2 * PAD, 0.45, t, size=12, color=color, line=1.3)

def card(s, x, y, w, h, kind="tile"):
    fill = {"tile": TILE, "sage": SAGE, "dark": JET, "hot": FLAME, "paper2": PAPER2}[kind]
    rect(s, x, y, w, h, fill, radius=0.16)
    return fill

def card_text(s, x, y, w, h, kind, title, body, kn=None, badge=None, logos=None, title_size=17, body_size=13):
    """Карточка: якорь сверху (бейдж-иконка или ряд логотипов), текст снизу."""
    card(s, x, y, w, h, kind)
    dark = kind in ("dark",); hot = kind == "hot"
    tcol = PAPER_ON_DARK if dark else INK
    bcol = C("C9C4B8") if dark else (C("3E2A1C") if hot else MUTED)
    kcol = FLAME if dark else (C("5A3416") if hot else FLAME_INK)
    px, py = x + 0.28, y + 0.28
    cy = py
    if kn:
        text(s, px, cy, w - 0.56, 0.25, kn, size=9, bold=True, color=kcol, caps=True); cy += 0.32
    if badge:
        d = 1.15
        oval(s, px, cy, d, (C("2A2822") if dark else (C("D9560F") if hot else PAPER2)))
        pic(s, badge, px + 0.26, cy + 0.26, d - 0.52, d - 0.52)
    if logos:
        lx = px; lh = 0.62
        for lp in logos:
            iw, ih = Image.open(lp).size; lw = min(1.25, lh * iw / ih)
            pic(s, lp, lx, cy, lw, lh); lx += lw + 0.18
    # текстовый блок прижат к низу; если якорь высокий — блок сдвигается под него
    import math
    anchor_bottom = cy + (1.15 if badge else (0.62 if logos else 0))
    avail = w - 0.56
    lines = sum(max(1, math.ceil(len(part) * title_size * 0.0095 / avail)) for part in title.split("\n"))
    th = 0.36 * lines * (title_size / 18) + 0.06
    body_h = 1.8
    block_h = th + 0.08 + body_h
    by = max(anchor_bottom + 0.18, y + h - 0.28 - block_h)
    text(s, px, by, w - 0.56, th, title, size=title_size, bold=True, color=tcol, line=1.08)
    text(s, px, by + th + 0.08, w - 0.56, min(body_h, y + h - 0.2 - (by + th + 0.08)), body, size=body_size, color=bcol, line=1.3)

def cards3(s, specs, y=2.45, h=4.15):
    gap = 0.2; w = (W - 2 * PAD - 2 * gap) / 3
    for i, sp in enumerate(specs):
        card_text(s, PAD + i * (w + gap), y, w, h, **sp)

# ── содержание ─────────────────────────────────────────────────────────────
def cover():
    s = new_slide(FLAME); top(s, "ПАС · ИТ-1 · 3 курс · 2026/27", hot=True)
    text(s, PAD, 2.1, 9, 1.8, "ГЕЙМДЕВ", size=96, bold=True, color=INK, line=1.0)
    text(s, PAD, 4.05, 8, 1.0, "Такт 2. Объекты конкуренции\nи бутылочные горлышки", size=22, bold=False, color=INK, line=1.2)
    s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(PAD), Inches(5.55), Inches(W - PAD), Inches(5.55)).line.color.rgb = C("8A4A28")
    text(s, PAD, 5.75, 6, 0.4, "Разбираем на Cyberpunk 2077 и CD Projekt", size=12, bold=True, color=INK)
    text(s, 7.2, 5.75, 5.5, 0.4, "8 этажей · 1 горлышко · прототип для правок", size=12, bold=True, color=INK, align=PP_ALIGN.RIGHT)

def section(num, title, sub):
    s = new_slide(JET); top(s, f"Раздел {num}", dark=True)
    text(s, W - PAD - 4.2, 0.35, 4.2, 2.6, num, size=150, bold=True, color=FLAME, align=PP_ALIGN.RIGHT, line=0.9)
    text(s, PAD, 3.4, 9, 2.0, title, size=64, bold=True, color=PAPER_ON_DARK, line=0.98, caps=True)
    text(s, PAD, 5.6, 9, 0.6, sub, size=14, color=C("B8B3A7"))

def slide_gamedev():
    s = new_slide(); top(s, "Определение"); h2(s, "Что такое геймдев")
    w1 = 7.2; w2 = W - 2 * PAD - w1 - 0.2
    card(s, PAD, 2.45, w1, 4.15, "tile")
    oval(s, PAD + 0.28, 2.73, 1.35, PAPER2); pic(s, ui_png("play"), PAD + 0.58, 3.03, 0.75, 0.75)
    text(s, PAD + 0.28, 4.45, w1 - 0.56, 2.0,
         [[ "Программные интерактивные системы, в которых пользователь ", ("добровольно действует в рамках искусственных правил", {"bold": True}),
            ", а система непрерывно отвечает на его действия, — с доведением продукта и его поддержки до аудитории ",
            ("через цифровые площадки дистрибуции", {"bold": True}), "."]], size=15, color=INK, line=1.3)
    card_text(s, PAD + w1 + 0.2, 2.45, w2, 4.15, "dark", "Не геймдев",
              "Веб-дизайн, тренажёры и обучающие симуляторы, классические приложения, розничная продажа игр, настольные игры.",
              badge=ui_png("no"))

def slide_object():
    s = new_slide(); top(s, "Понятие"); h2(s, "Объект конкуренции — то,\nобладание которым даёт преимущество")
    cards3(s, [
        dict(kind="tile", title="Право, стандарт,\nопыт, внимание", body="Это можно удерживать. А пока удерживаешь — другим сюда хода нет.", badge=ui_png("key")),
        dict(kind="sage", title="Но не деньги\nи не доля рынка", body="Это ресурс на входе и следствие на выходе. За них не дерутся — их получают.", badge=ui_png("nomoney")),
        dict(kind="dark", title="Проверка: сравни\nс этажом выше", body="Объект тот же — значит, это один уровень, а не два. Если объект нельзя измерить — мы его не называем.", badge=ui_png("compare")),
    ])

ROWS = [
 ("01", "Интеллектуальная\nсобственность", "Право на мир и героев: лицензии на франшизы", ["nintendo", "sega", "cdprojekt"]),
 ("02", "Концепт", "Решение о финансировании — greenlight", ["valve", "rockstargames"]),
 ("03", "Движки\nи инструменты", "Выбор студий: кто станет стандартом", ["unrealengine", "unity", "godotengine"]),
 ("04", "Компоненты\nи аутсорс", "Контракты разработчиков и издателей", ["virtuos", "room8", "keywords", "jali"]),
 ("05", "Разработка", "Выбор игроком игры — какую играть", ["cdprojekt", "rockstargames", "riotgames", "valve"]),
 ("06", "Издание\nи маркетинг", "Права на издание и релиз", ["sony", "ubisoft", "ea", "squareenix"]),
 ("07", "Платформы\nи железо", "Установленная база и эксклюзивы", ["playstation", "xbox", "nintendoswitch", "apple", "android"]),
 ("08", "Витрины\nи дистрибуция", "Аккаунт и библиотека игрока; каталог разработчиков", ["steam", "playstation", "xbox", "nintendoswitch", "appstore", "googleplay"]),
]
def slide_map():
    s = new_slide(); top(s, "Схема · разметка карты")
    y0 = 1.05; text(s, PAD + 0.55, y0, 2.4, 0.25, "УРОВЕНЬ", size=8.5, bold=True, color=MUTED)
    text(s, PAD + 3.05, y0, 4, 0.25, "ОБЪЕКТ КОНКУРЕНЦИИ", size=8.5, bold=True, color=MUTED)
    text(s, W - PAD - 3.6, y0, 3.6, 0.25, "КТО СТОИТ", size=8.5, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
    rh = 0.66; gap = 0.07; y = y0 + 0.32
    for i, (no, nm, oj, logos) in enumerate(ROWS):
        neck = (no == "08")
        rect(s, PAD, y, W - 2 * PAD, rh, JET if neck else TILE, radius=0.12)
        col = FLAME if neck else MUTED
        text(s, PAD + 0.16, y, 0.4, rh, no, size=14, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
        text(s, PAD + 0.62, y, 2.4, rh, nm, size=13.5, bold=True, color=(FLAME if neck else INK), anchor=MSO_ANCHOR.MIDDLE, line=1.05)
        text(s, PAD + 3.05, y, 5.2, rh, oj, size=13.5, bold=True, color=(C("FF9A52") if neck else FLAME_INK), anchor=MSO_ANCHOR.MIDDLE)
        lh = 0.4; lx = W - PAD - 0.2
        for slug in reversed(logos):
            p = ico(slug, on_dark=neck); iw, ih = Image.open(p).size; lw = min(0.95, lh * iw / ih)
            lx -= lw; pic(s, p, lx, y + (rh - lh) / 2, lw, lh); lx -= 0.16
        y += rh + gap

def floor(no, name, obj, title, specs, cp):
    s = new_slide(); top(s, f"Этаж {no} · {name}"); h2(s, title)
    cards3(s, specs, y=2.35, h=4.1)
    text(s, PAD, 6.6, W - 2 * PAD, 0.6, [[("На примере Cyberpunk 2077: ", {"bold": True, "color": FLAME_INK}), cp]], size=12, color=MUTED, line=1.3)

def floors():
    floor("01", "интеллектуальная собственность", "право на мир и героев",
          "Интеллектуальная собственность:\nправо на мир и героев",
          [dict(kind="tile", title="Если ты не Nintendo —\nне сделаешь Mario", body="Право исключительное: у одного есть, у остальных нет. Mario 1985 года всё ещё закрывает этаж.", logos=[ico("nintendo")]),
           dict(kind="tile", title="Барьер не денежный,\nа юридический", body="Купить нельзя, если владелец не продаёт. Но свою ИС можно создать с нуля — Half-Life, Minecraft — значит, не горлышко.", badge=ui_png("key")),
           dict(kind="sage", title="За что конкурируют", body="За лицензии на франшизы — Marvel, Star Wars, Warhammer. Кто получил права, тот и делает игру по миру.", badge=ui_png("shield"))],
          "CD Projekt не придумала Найт-Сити — взяла права на настольную Cyberpunk 2020 у Майка Пондсмита.")
    floor("02", "концепт", "решение о финансировании",
          "Концепт:\nрешение о финансировании (greenlight)",
          [dict(kind="tile", title="Что такое концепт", body="Документ, по которому считают бюджет, сроки и риск. «Хотим RPG про киберпанк» — это ещё не он.", badge=ui_png("doc")),
           dict(kind="sage", title="За что конкурируют", body="За решение дать денег. Внутри Rockstar концепты конкурируют за greenlight между собой; инди — за издателя или платформенный грант снаружи.", badge=ui_png("compare")),
           dict(kind="dark", title="Риск разный —\nобъект один", body="Rockstar ставит миллиарды и 13–15 лет, инди — год и тысячи долларов. Не «низкий» и «высокий» риск, а разный масштаб одной ставки на greenlight.", badge=ui_png("layers"))],
          "концепт Пондсмита прошёл greenlight в CD Projekt после успеха «Ведьмака 3» — на ресурсах и репутации студии.")
    floor("03", "движки и инструменты", "выбор студий",
          "Движки и инструменты:\nвыбор студий",
          [dict(kind="tile", title="Два движка — у 72% студий", body="Unreal — 42%, Unity — 30% по опросу GDC 2026. Godot — 8–10% релизов в Steam. Кого выбирают студии, тот становится стандартом.", logos=[ico("unrealengine"), ico("unity"), ico("godotengine")]),
           dict(kind="sage", title="Стандарт — следствие", body="Под движок нанимают людей, строят пайплайн, покупают ассеты. Сменить — значит переучить студию. Но можно: Valve ушла с Quake на GoldSrc.", badge=ui_png("layers")),
           dict(kind="dark", title="Свой движок: RAGE, Source", body="Rockstar и Valve держат свои. Закрытость ценой выбора подрядчиков: этаж 03 управляет этажом 04.", logos=[ico("rockstargames", True), ico("valve", True)])],
          "сделан на своём REDengine; в 2022-м CD Projekt объявила переход на Unreal Engine 5 — свой движок оказался заменяемым.")
    floor("04", "компоненты и аутсорс", "контракты",
          "Компоненты и аутсорс:\nконтракты разработчиков и издателей",
          [dict(kind="tile", title="Сотни студий\nс узким профилем", body="Virtuos и Keywords — производство, Room 8 — арт, JALI — лицевая анимация. Конкурируют за контракты; оружие — опыт: опыт → качество → цена → срок.", logos=[ico("virtuos"), ico("room8"), ico("keywords"), ico("jali")]),
           dict(kind="sage", title="Отбор формальный", body="Титры похожих игр, фильтр по профилю, арт-тест, пилот на реальном контенте — и только потом контракт.", badge=ui_png("compare")),
           dict(kind="dark", title="Контрмодель: Rockstar", body="Свои студии и свой движок. Внешнего арт-аутсорса почти нет — ценой закрытости.", logos=[ico("rockstargames", True)])],
          "собран из компонентов внешних команд — механики, лицевая анимация, звук, QA — всё разными студиями.")
    floor("05", "разработка", "выбор игроком игры",
          "Разработка:\nвыбор игроком игры",
          [dict(kind="tile", title="Кто делает саму игру", body="CD Projekt RED, Rockstar North, Riot, Valve, Insomniac. Самый населённый этаж — тысячи студий.", logos=[ico("cdprojekt"), ico("rockstargames"), ico("riotgames"), ico("valve")]),
           dict(kind="sage", title="Какую играть —\nне где покупать", body="Разработчик конкурирует за выбор игры игроком. Витрина — за то, где он её купит. Два разных объекта.", badge=ui_png("eye")),
           dict(kind="dark", title="Заменяем", body="Студий много, профили пересекаются. Игра переживает смену разработчика — Risk of Rain пережила.", badge=ui_png("layers"))],
          "разработчик — CD Projekt RED, около 500 человек в пике производства.")
    floor("06", "издание и маркетинг", "права на издание",
          "Издание:\nправа на издание и релиз",
          [dict(kind="tile", title="Издатель ≠ разработчик", body="«Человека-паука» 2018 года сделала Insomniac — тогда независимая студия. Sony была издателем и купила её через год за $229 млн.", kn="Издатель", logos=[ico("sony"), ico("ubisoft"), ico("ea"), ico("squareenix")]),
           dict(kind="sage", title="Холдинг — доли\nв тех, кто делает", body="Tencent взяла 93% Riot в 2011-м и весь остаток к 2015-му. Игр не делает сама.", kn="Холдинг", badge=ui_png("layers")),
           dict(kind="dark", title="Три роли,\nтри объекта", body="На старой карте Tencent и Sony стояли на «разработке». Разработчик, издатель и холдинг конкурируют за разное.", kn="Ошибка старой карты", badge=ui_png("compare"))],
          "у CD Projekt издатель и разработчик — разные юрлица одной группы.")
    floor("07", "платформы и железо", "установленная база",
          "Платформы и железо:\nустановленная база и эксклюзивы",
          [dict(kind="tile", title="Консоли", body="Закрытые экосистемы со своими правилами и сертификацией. Конкурируют за игроков с устройством и за эксклюзивы.", logos=[ico("playstation"), ico("xbox"), ico("nintendoswitch")]),
           dict(kind="sage", title="Мобайл", body="Самая большая аудитория — и свои ворота у каждой системы.", logos=[ico("apple"), ico("android")]),
           dict(kind="dark", title="PC — открытая", body="Windows, Mac, Linux. Вышел на Windows и не вышел на Mac — потерял пару процентов. Платформы — данность: под них делают.", badge=ui_png("gate"))],
          "провал запуска случился на PS4 и Xbox One — платформа определила судьбу релиза и снятие с PS Store.")

def slide_showcase():
    s = new_slide(); top(s, "Этаж 08 · витрины и дистрибуция")
    h2(s, "Витрины и дистрибуция:\nаккаунт и библиотека игрока")
    cards3(s, [
        dict(kind="tile", title="Что такое витрина", body="Платформа, которая владеет аккаунтом игрока: библиотекой, платежами, рекомендациями. Steam, PlayStation Store, Xbox, eShop, App Store, Google Play.",
             logos=[ico("steam"), ico("playstation"), ico("xbox"), ico("appstore"), ico("googleplay")]),
        dict(kind="sage", title="За что конкурируют витрины", body="За аккаунт с библиотекой — он создаёт издержки переключения. За эксклюзивы: CS, Dota, Half-Life только в Steam. За каталог разработчиков — комиссией 30% против 12%.", badge=ui_png("key")),
        dict(kind="hot", title="Почему горлышко", body="Готовая игра не доходит до игрока иначе как через одну из немногих таких платформ, и каждая контролирует свою экосистему. Внимание и привычка — следствие, а не объект.", badge=ui_png("gate")),
    ], y=2.35, h=4.1)
    text(s, PAD, 6.6, W - 2 * PAD, 0.5, [[("На примере Cyberpunk 2077: ", {"bold": True, "color": FLAME_INK}), "снята с PlayStation Store 18.12.2020 с полным возвратом денег, вернулась 21.06.2021."]], size=12, color=MUTED)

def slide_chain():
    s = new_slide(); top(s, "Этаж 08 · цепочка удержания"); h2(s, "Как время игрока превращается\nв издержки переключения")
    steps = [("Время\nна платформе", "ресурс, не объект"), ("Покупки", "витрина берёт комиссию"), ("Библиотека", "не переносится"),
             ("Друзья, достижения,\nотзывы", "живут только здесь"), ("Издержки\nпереключения", "уйти = бросить всё"), ("LTV игрока", "деньги дальше по цепочке")]
    n = len(steps); gap = 0.22; w = (W - 2 * PAD - gap * (n - 1)) / n; y = 2.7; h = 2.2
    for i, (t, sub) in enumerate(steps):
        x = PAD + i * (w + gap); last = i == n - 1
        rect(s, x, y, w, h, FLAME if last else TILE, radius=0.14)
        text(s, x + 0.18, y + 0.25, w - 0.36, 1.1, t, size=14, bold=True, color=INK, line=1.08)
        text(s, x + 0.18, y + h - 0.75, w - 0.36, 0.55, sub, size=11, color=(C("3E2A1C") if last else MUTED))
        if not last:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + w + 0.02), Inches(y + h / 2 - 0.1), Inches(gap - 0.04), Inches(0.2))
            a.fill.solid(); a.fill.fore_color.rgb = FLAME; a.line.fill.background()
    text(s, PAD, 5.35, W - 2 * PAD, 1.2,
         [["Отзывы в Steam — артефакт доверия, которого нет у Epic в таком виде: чем больше игроков, тем качественнее отзывы, тем выше готовность покупать и рисковать. Достижения не переносятся, друзья для онлайна — только через платформу. Всё это и есть удержание."]],
         size=13, color=MUTED, line=1.35)
    note(s, "Правило: время — ресурс, который конвертируется в удержание. Объектом мы называем то, что накапливается и удерживается, — библиотеку и аккаунт.")

def slide_pipeline():
    s = new_slide(); top(s, "Этаж 08 · зависимость разработчика"); h2(s, "Что нужно разработчику,\nчтобы игра появилась на витрине")
    steps = [("Сертификация", "техтребования каждой платформы; блокер"), ("Комиссия", "Steam 30% · Epic 12%\nApp Store, Google Play 30% (15% малому бизнесу)\nконсоли ≈30%, условия под NDA"),
             ("Отдельный билд", "под каждую платформу; порт — отдельный бюджет"), ("Возрастной рейтинг", "ESRB, PEGI, USK; без него — не на консоль"), ("Правила контента", "модерация витрины; отказ = нет релиза")]
    n = len(steps); gap = 0.2; w = (W - 2 * PAD - gap * (n - 1)) / n; y = 2.6; h = 3.0
    for i, (t, sub) in enumerate(steps):
        x = PAD + i * (w + gap)
        rect(s, x, y, w, h, TILE, radius=0.14)
        text(s, x + 0.2, y + 0.22, 0.6, 0.4, f"{i+1}", size=20, bold=True, color=FLAME)
        text(s, x + 0.2, y + 0.75, w - 0.4, 0.55, t, size=15, bold=True, color=INK, line=1.08)
        text(s, x + 0.2, y + 1.35, w - 0.4, h - 1.55, sub, size=11.5, color=MUTED, line=1.3)
    note(s, "Каждый шаг — отдельные ворота со своей ценой, и условия диктует витрина. Быть на всех витринах дорого — и это работает на тезис о горлышке.", y=5.85)
    text(s, PAD, 6.5, W - 2 * PAD, 0.5, "Роялти Unreal — 5% после $1 млн, но 0% при продаже в Epic Games Store: витрина Epic субсидируется движком Epic.", size=12, color=FLAME_INK, bold=True)

def slide_measure():
    s = new_slide(JET); top(s, "Этаж 08 · измеримость объекта", dark=True)
    h2(s, "Если объект нельзя измерить —\nмы его не называем", color=PAPER_ON_DARK)
    figs = [("147 млн", "ежемесячных активных игроков Steam, 2025"), ("42 млн", "рекордный одновременный онлайн Steam, март 2026"), ("295 млн", "аккаунтов в Epic Games Store при 8–10% доли")]
    w = (W - 2 * PAD) / 3
    for i, (n, l) in enumerate(figs):
        x = PAD + i * w
        text(s, x, 2.7, w - 0.4, 1.0, n, size=54, bold=True, color=FLAME, line=0.95)
        text(s, x, 3.75, w - 0.5, 0.8, l, size=12.5, color=C("B8B3A7"), line=1.3)
        if i: s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x - 0.2), Inches(2.75), Inches(x - 0.2), Inches(4.4)).line.color.rgb = C("3A372F")
    for i, (t, b) in enumerate([("Что измеряем", "MAU и одновременный онлайн — размер аккаунтной базы; размер библиотек и вишлистов — накопленное удержание; число и объём отзывов — доверие."),
                                 ("Чего пока нет", "Публичного размера библиотек и вишлистов по платформам. Пока не измерили — не заявляем как объект; это открытая задача.")]):
        x = PAD + i * ((W - 2 * PAD) / 2 + 0.1); w2 = (W - 2 * PAD) / 2 - 0.1
        rect(s, x, 4.85, w2, 1.65, C("232119"), radius=0.14)
        text(s, x + 0.28, 5.05, w2 - 0.56, 0.4, t, size=15, bold=True, color=PAPER_ON_DARK)
        text(s, x + 0.28, 5.5, w2 - 0.56, 0.95, b, size=12, color=C("B8B3A7"), line=1.3)

def slide_multi():
    s = new_slide(); top(s, "Кто стоит на нескольких этажах"); h2(s, "Позиция на одном этаже —\nрычаг на другом")
    groups = [("Valve", [("03 · Движок", "valve", "Source"), ("05 · Разработка", "valve", "Half-Life, CS, Dota"), ("08 · Витрина", "steam", "Steam — 75% PC")]),
              ("Epic Games", [("03 · Движок", "unrealengine", "Unreal — 42% студий"), ("05 · Разработка", "epicgames", "Fortnite"), ("08 · Витрина", "epicgames", "EGS — 8–10%")])]
    gw = (W - 2 * PAD - 0.25) / 2
    for gi, (name, items) in enumerate(groups):
        gx = PAD + gi * (gw + 0.25); rect(s, gx, 2.45, gw, 2.75, PAPER2, radius=0.16)
        text(s, gx + 0.25, 2.62, gw - 0.5, 0.35, name, size=15, bold=True, color=INK)
        tw = (gw - 0.5 - 0.3) / 3
        for ti, (lab, slug, sub) in enumerate(items):
            tx = gx + 0.25 + ti * (tw + 0.15); rect(s, tx, 3.05, tw, 2.0, TILE, radius=0.13)
            pic(s, ico(slug), tx + tw / 2 - 0.42, 3.22, 0.84, 0.84)
            text(s, tx + 0.12, 4.15, tw - 0.24, 0.3, lab, size=9, bold=True, color=FLAME_INK, align=PP_ALIGN.CENTER, caps=True)
            text(s, tx + 0.12, 4.45, tw - 0.24, 0.5, sub, size=10.5, color=MUTED, align=PP_ALIGN.CENTER, line=1.2)
    text(s, PAD, 5.4, W - 2 * PAD, 1.2,
         [["Доля витрины у Epic — 8–10%, но 42% студий делают игры на его движке. Сила на этаже 03 финансирует борьбу на этаже 08: эксклюзивы, бесплатные раздачи, комиссия 12%, 0% роялти в своём магазине. Valorant сделан Riot на Unreal — Epic получает роялти с чужого хита. Доля на одном этаже — не мера силы компании."]],
         size=13, color=MUTED, line=1.35)
    note(s, "Гипотеза Владимира к проверке: ценность магазина для Epic — не выручка, а доверие игроков и разработчиков к экосистеме Unreal.")

def slide_vs():
    s = new_slide(); top(s, "Проверка · объекты различаются"); h2(s, "Объекты соседних этажей\nне совпадают")
    pairs = [("03 · Движки", "выбор студий: раз выбрали — живут годами", "04 · Аутсорс", "контракты: торгуются заново на каждом проекте"),
             ("05 · Разработка", "какую игру выберет игрок", "06 · Издание", "права на издание и релиз — не про игрока"),
             ("07 · Платформа", "установленная база: у кого есть устройство", "08 · Витрина", "аккаунт и библиотека: где игрок живёт")]
    y = 2.5; h = 1.25
    for a, at, b, bt in pairs:
        rect(s, PAD, y, W - 2 * PAD, h, TILE, radius=0.14)
        half = (W - 2 * PAD - 0.8) / 2
        for k, (lab, t) in enumerate(((a, at), (b, bt))):
            x = PAD + 0.3 + k * (half + 0.8)
            text(s, x, y + 0.22, half - 0.2, 0.3, lab, size=9.5, bold=True, color=MUTED, caps=True)
            text(s, x, y + 0.55, half - 0.2, 0.6, t, size=14, bold=True, color=INK, line=1.15)
        text(s, PAD + 0.3 + half, y, 0.8, h, "≠", size=28, bold=True, color=FLAME, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += h + 0.18

def slide_neck_def():
    s = new_slide(); top(s, "Понятие · бутылочное горлышко"); h2(s, "Уровень, где мало игроков,\nа зависят от них все")
    lw = 7.4; specs = [("Признак 1", "Мало игроков", "Считаем тех, к кому реально можно пойти. Сотня студий — много, несколько витрин — мало.", "tile"),
                       ("Признак 2", "Нет альтернативы", "Можно заменить — больно, дорого, но можно — это не горлышко.", "tile"),
                       ("Что удерживает", "Барьер входа", "Не деньги, а то, что деньгами не покупается: сетевой эффект, библиотека, привычка.", "hot")]
    y = 2.5; ch = 1.28
    for kn, t, b, kind in specs:
        card(s, PAD, y, lw, ch, kind); dark = kind == "hot"
        text(s, PAD + 0.28, y + 0.2, lw - 0.56, 0.25, kn, size=9, bold=True, color=(C("5A3416") if dark else FLAME_INK), caps=True)
        text(s, PAD + 0.28, y + 0.47, lw - 0.56, 0.35, t, size=16, bold=True, color=INK)
        text(s, PAD + 0.28, y + 0.82, lw - 0.56, 0.45, b, size=11.5, color=(C("3E2A1C") if dark else MUTED), line=1.25)
        y += ch + 0.15
    # воронка
    fx = PAD + lw + 0.25; fw = W - PAD - fx; rect(s, fx, 2.5, fw, 4.15, TILE, radius=0.16)
    text(s, fx, 2.75, fw, 0.3, "ЭТАЖИ 01–07", size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    cone = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(fx + 0.5), Inches(3.1), Inches(fw - 1.0), Inches(2.1))
    cone.rotation = 180; cone.adjustments[0] = 0.42; cone.fill.solid(); cone.fill.fore_color.rgb = PAPER2; cone.line.fill.background()
    for k in range(1, 7):
        yy = 3.1 + 2.1 * k / 7; ins = (fw - 1.0) * 0.42 * k / 7
        s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(fx + 0.5 + ins), Inches(yy), Inches(fx + fw - 0.5 - ins), Inches(yy)).line.color.rgb = TILE
    nw = (fw - 1.0) * 0.16 + 0.1; nx = fx + fw / 2 - nw / 2
    rect(s, nx, 5.2, nw, 0.7, FLAME, radius=0)
    text(s, nx, 5.22, nw, 0.4, "08", size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, nx, 5.58, nw, 0.3, "ВИТРИНА", size=7, bold=True, color=C("3E2A1C"), align=PP_ALIGN.CENTER)
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(fx + fw / 2 - 0.12), Inches(5.95), Inches(0.24), Inches(0.4)); a.fill.solid(); a.fill.fore_color.rgb = FLAME; a.line.fill.background()
    text(s, fx, 6.35, fw, 0.3, "ИГРОК", size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)

def slide_failtest():
    s = new_slide(); top(s, "Цепочка зависимостей · тест отказом"); h2(s, "Тест отказом\nпо четырём уровням")
    items = [("Откажет движок", "Берёшь другой. Valve ушла с Quake и собрала GoldSrc; CD Projekt уходит с REDengine на Unreal.", "Альтернатива есть", False),
             ("Откажет подрядчик", "Меняешь. Рынок аутсорса — около $9 млрд и сотни студий.", "Альтернатива есть", False),
             ("Откажет издатель", "Ищешь другого или идёшь без него. Игрок этого даже не заметит.", "Есть, но дорогая", False),
             ("Откажет витрина", "Игры не существует. Игрок физически не может её найти и купить.", "Альтернативы нет", True)]
    gap = 0.2; w = (W - 2 * PAD - 3 * gap) / 4; y = 2.5; h = 3.8
    for i, (t, b, v, dead) in enumerate(items):
        x = PAD + i * (w + gap); rect(s, x, y, w, h, JET if dead else TILE, radius=0.16)
        text(s, x + 0.25, y + 0.28, w - 0.5, 0.7, t, size=16, bold=True, color=(FLAME if dead else INK), line=1.1)
        text(s, x + 0.25, y + 1.05, w - 0.5, 1.9, b, size=12.5, color=(C("C9C4B8") if dead else MUTED), line=1.3)
        text(s, x + 0.25, y + h - 0.7, w - 0.5, 0.45, v, size=13, bold=True, color=(FLAME if dead else INK))
    note(s, "Три уровня из четырёх заменяемы. Жёсткая зависимость ровно одна.")

def slide_class():
    s = new_slide(JET); top(s, "Горлышко · обоснование", dark=True)
    text(s, PAD, 1.05, W - 2 * PAD, 1.25, [["Горлышко — ", ("уровень витрин", {"color": FLAME}), ",\nа не отдельная компания"]], size=36, bold=True, color=PAPER_ON_DARK, line=1.05)
    figs = [("75%", "выручки PC-дистрибуции у Steam — это лишь PC-проекция горлышка"), ("2–3%", "у GOG — третьей витрины PC. Всё, что не Steam и не Epic, — крошки"), ("8–10%", "у Epic, который раздаёт игры бесплатно не первый год")]
    w = (W - 2 * PAD) / 3
    for i, (n, l) in enumerate(figs):
        x = PAD + i * w
        text(s, x, 2.75, w - 0.4, 1.0, n, size=58, bold=True, color=FLAME, line=0.95)
        text(s, x, 3.8, w - 0.5, 0.8, l, size=12, color=C("B8B3A7"), line=1.3)
    for i, (t, b) in enumerate([("Что мешает войти", "Купленная библиотека не переносится. Уйти — значит бросить всё, за что заплачено."),
                                 ("Контраргумент", "Steam — только PC. Поэтому горлышко не компания, а уровень: в каждой экосистеме своя витрина.")]):
        x = PAD + i * ((W - 2 * PAD) / 2 + 0.1); w2 = (W - 2 * PAD) / 2 - 0.1
        rect(s, x, 4.9, w2, 1.6, C("232119"), radius=0.14)
        text(s, x + 0.28, 5.1, w2 - 0.56, 0.4, t, size=15, bold=True, color=PAPER_ON_DARK)
        text(s, x + 0.28, 5.55, w2 - 0.56, 0.9, b, size=12, color=C("B8B3A7"), line=1.3)

def slide_cases():
    s = new_slide(); top(s, "Горлышко · проверенные случаи"); h2(s, "Три случая отключения\nигры витриной")
    cases = [("playstation", "Cyberpunk 2077", "6 месяцев вне PS Store", "Sony убрала игру 18.12.2020 и вернула деньги всем. Вернулась 21.06.2021. Акции CD Projekt — минус 20% за день."),
             ("apple", "Fortnite", "почти 5 лет вне App Store", "Apple удалила игру в августе 2020-го. Вернулась в американский App Store в мае 2025-го. Пять лет игры не существовало для iOS."),
             ("vk", "Atomic Heart", "только VK Play в РФ и СНГ", "На релизе 21.02.2023 страница в Steam была недоступна в регионе. В российском Steam — только в августе 2026-го.")]
    cards3(s, [dict(kind="tile", title=t, body=b, kn=sub, badge=ico(sl)) for sl, t, sub, b in cases], y=2.45, h=3.95)
    note(s, "Игра была готова и куплена. Отказал только уровень витрины — и этого хватило.")

def slide_hidden():
    s = new_slide(); top(s, "Скрытые горлышки"); h2(s, "Скрытые горлышки:\nсертификация, издатель, мощности")
    cards3(s, [dict(kind="tile", title="Сертификация\nи возрастной рейтинг", body="Не прошёл ESRB или PEGI — на консоль не вышел. Обойти нельзя.", badge=ui_png("shield")),
               dict(kind="tile", title="Нет издателя —\nнет игры", body="Команда сильная, права есть, а издателя нет. Для инди горлышко — шестой этаж, а не восьмой.", badge=ui_png("no")),
               dict(kind="tile", title="Мощности подрядчиков", body="Cyberpunk сорвал не аутсорс, а скоуп и старые консоли. Но конвейер студий — признак нехватки.", badge=ui_png("clock"))], y=2.45, h=3.95)
    note(s, "Горлышко ищется не по выручке этажа, а по отсутствию альтернативы у тех, кто зависит.")

def slide_open():
    s = new_slide(FLAME); top(s, "Открытый вопрос · выносим в зал", hot=True)
    h2(s, "Открытый вопрос:\nгде проходит граница рынка")
    cw = 4.6; gap = 0.25; x0 = (W - 2 * cw - gap) / 2; y = 3.3; h = 1.9
    for i, (kn, t, b) in enumerate([("Рынок = PC", "Горлышко — Steam", "75% выручки, альтернатив практически нет."),
                                     ("Рынок = гейминг целиком", "Горлышек несколько", "По одному на экосистему — и каждое монополист в своей.")]):
        x = x0 + i * (cw + gap); rect(s, x, y, cw, h, TILE, radius=0.16)
        text(s, x + 0.28, y + 0.28, cw - 0.56, 0.25, kn, size=9, bold=True, color=FLAME_INK, caps=True)
        text(s, x + 0.28, y + 0.6, cw - 0.56, 0.45, t, size=18, bold=True, color=INK)
        text(s, x + 0.28, y + 1.1, cw - 0.56, 0.7, b, size=13, color=MUTED, line=1.3)

SRC = [
 ("Unreal 42%, Unity 30% — основной движок студий (опрос GDC 2026); Godot 8–10% релизов Steam", "этаж 03", "shattered.io", "https://shattered.io/unreal-engine-6-unity-market-share-2026"),
 ("Steam 74–75% PC-дистрибуции; Epic 8–10%; GOG 2–3%; 147 млн MAU; пик 42 млн", "этаж 08, горлышко", "levvvel.com", "https://levvvel.com/statistics/pc-gaming"),
 ("Комиссии: Steam 30%, Epic 12%; с 06.2025 у Epic 0% на первый $1 млн в год", "пайплайн витрины", "shattered.io", "https://shattered.io/steam-vs-epic-games-store/"),
 ("Epic Games Store: 295 млн аккаунтов, $1,09 млрд трат за год", "измеримость", "shattered.io", "https://shattered.io/epic-games-store-revenue-2026/"),
 ("Роялти Unreal 5% после $1 млн; 0% при продаже в Epic Games Store", "рычаг Epic", "unrealengine.com", "https://www.unrealengine.com/license"),
 ("App Store 30%, 15% для малого бизнеса", "пайплайн витрины", "developer.apple.com", "https://developer.apple.com/app-store/small-business-program/"),
 ("Google Play 30%, 15% на первый $1 млн", "пайплайн витрины", "support.google.com", "https://support.google.com/googleplay/android-developer/answer/112622"),
 ("CD Projekt: переход с REDengine на Unreal Engine 5, март 2022", "этаж 03, тест отказом", "cdprojekt.com", "https://www.cdprojekt.com/en/media/news/new-witcher-saga-announced-cd-projekt-red-begins-development-on-unreal-engine-5-as-part-of-a-strategic-partnership-with-epic-games/"),
 ("Аутсорс ≈ $9 млрд, 2025 (другие отчёты: $4–9 млрд)", "тест отказом", "verifiedmarketreports.com", "https://www.verifiedmarketreports.com/product/game-outsourcing-service-market/"),
 ("Sony купила Insomniac за $229 млн, август 2019", "этаж 06", "gameworldobserver.com", "https://gameworldobserver.com/2020/02/11/sony-paid-229m-insomniac-games"),
 ("Tencent: 93% Riot в 2011, 100% к концу 2015", "этаж 06", "techcrunch.com", "https://techcrunch.com/2015/12/16/tencent-takes-full-control-of-league-of-legends-creator-riot-games/"),
 ("Cyberpunk вне PS Store 18.12.2020 — 21.06.2021; акции CD Projekt −20%", "кейсы", "cnbc.com", "https://www.cnbc.com/2020/12/18/sony-pulls-cyberpunk-2077-from-playstation-store-after-backlash.html"),
 ("Fortnite вне App Store: август 2020 — май 2025", "кейсы", "cnbc.com", "https://www.cnbc.com/2025/05/20/apple-fortnite-app-store-epic-games.html"),
 ("Atomic Heart только в VK Play с 21.02.2023; в Steam РФ — 08.2026", "кейсы", "habr.com", "https://habr.com/ru/news/687214/"),
 ("Super Mario Bros. — 1985", "этаж 01", "nintendo.com", "https://www.nintendo.com/us/store/products/super-mario-bros-nes-nintendo-switch-online/"),
]
def slide_sources():
    s = new_slide(); top(s, "Источники всех чисел"); h2(s, "Откуда взяты числа", size=30)
    rows = len(SRC) + 1; tb = s.shapes.add_table(rows, 3, Inches(PAD), Inches(1.95), Inches(W - 2 * PAD), Inches(4.9)).table
    tb.columns[0].width = Inches(6.6); tb.columns[1].width = Inches(2.2); tb.columns[2].width = Inches(W - 2 * PAD - 8.8)
    def cell(r, c, t, size=10, bold=False, color=INK, url=None):
        ce = tb.cell(r, c); ce.fill.solid(); ce.fill.fore_color.rgb = TILE if r else PAPER2
        ce.margin_left = ce.margin_right = Inches(0.1); ce.margin_top = ce.margin_bottom = Inches(0.04)
        tf = ce.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; run = p.add_run(); run.text = t
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
        if url: run.hyperlink.address = url
    for c, hdr in enumerate(("ЧИСЛО", "ГДЕ В ДЕКЕ", "ИСТОЧНИК")): cell(0, c, hdr, 8.5, True, MUTED)
    for r, (n, w, dom, url) in enumerate(SRC, 1):
        cell(r, 0, n, 10, True); cell(r, 1, w, 10, False, MUTED); cell(r, 2, dom, 10, True, FLAME_INK, url)
    note(s, "Консольные комиссии (≈30%) платформы публично не раскрывают — указаны как оценка, без ссылки.", y=6.95)

def slide_terms():
    s = new_slide(); top(s, "Понятия такта 2"); h2(s, "Понятия, которыми пользовались", size=30)
    gl = [("Конкуренция", "Борьба за то, чего на всех не хватает. Всегда за что-то конкретное."), ("Барьер входа", "Что мешает новому игроку занять уровень. Чем выше — тем дольше держится горлышко."),
          ("Объект конкуренции", "То, обладание которым даёт преимущество. На каждом этаже свой. Если нельзя измерить — не называем."), ("Зависимость в цепочке", "Насколько жёстко уровень нуждается в соседнем. Меряется отказом."),
          ("Бутылочное горлышко", "Уровень, где мало игроков, а зависят от них все."), ("Уровень", "Место, где есть свой объект конкуренции. Не шаг процесса.")]
    w = (W - 2 * PAD - 0.2) / 2; h = 1.35
    for i, (t, b) in enumerate(gl):
        x = PAD + (i % 2) * (w + 0.2); y = 2.0 + (i // 2) * (h + 0.15)
        rect(s, x, y, w, h, TILE, radius=0.14)
        text(s, x + 0.28, y + 0.25, w - 0.56, 0.4, t, size=15, bold=True, color=INK)
        text(s, x + 0.28, y + 0.68, w - 0.56, 0.6, b, size=12, color=MUTED, line=1.3)

# ── порядок ────────────────────────────────────────────────────────────────
cover(); slide_gamedev(); slide_object(); section("01", "Карта\nуровней", "Восемь этажей, на каждом свой объект конкуренции.")
slide_map(); floors(); slide_showcase(); slide_chain(); slide_pipeline(); slide_measure(); slide_multi(); slide_vs()
section("02", "Бутылочное\nгорлышко", "Ищем его тестом отказом, а не по деньгам на этаже.")
slide_neck_def(); slide_failtest(); slide_class(); slide_cases(); slide_hidden(); slide_open(); slide_sources(); slide_terms()

# проставляем общее число слайдов в счётчики
N = len(prs.slides)
for sl in prs.slides:
    for sh in sl.shapes:
        if sh.has_text_frame and "{N}" in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs: r.text = r.text.replace("{N}", str(N))
out = ROOT / "takt2-gamedev.pptx"; prs.save(out)
print("сохранено:", out.name, "| слайдов:", N, "| размер:", out.stat().st_size // 1024, "КБ")
